import asyncio
import json
import logging
import os
import shutil
import threading
import numpy as np
import lightrag.api.server_util as server_util
from docx import Document
from pptx import Presentation
from io import BytesIO
from PyPDF2 import PdfReader
from contextlib import asynccontextmanager
from pathlib import Path
from typing import List, Union, Dict, Optional
import aiofiles
from ascii_colors import trace_exception
from docling.document_converter import DocumentConverter
import fastapi as fastapi
import lightrag.llm.azure_openai as az_openai
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from lightrag import LightRAG, QueryParam
from lightrag.api import __api_version__
from lightrag.types import GPTKeywordExtractionFormat
from lightrag.utils import EmbeddingFunc

# Global progress tracker
scan_progress: Dict = {
    "is_scanning": False,
    "current_file": "",
    "indexed_count": 0,
    "total_files": 0,
    "progress": 0,
}

# Lock for thread-safe operations
progress_lock = threading.Lock()


def create_app(args, rag_storage_config):
    # Add SSL validation
    if args.ssl:
        if not args.ssl_certfile or not args.ssl_keyfile:
            raise Exception(
                "SSL certificate and key files must be provided when SSL is enabled"
            )
        if not os.path.exists(args.ssl_certfile):
            raise Exception(f"SSL certificate file not found: {args.ssl_certfile}")
        if not os.path.exists(args.ssl_keyfile):
            raise Exception(f"SSL key file not found: {args.ssl_keyfile}")

    # Setup logging
    logging.basicConfig(
        format="%(levelname)s:%(message)s", level=getattr(logging, args.log_level)
    )

    # Check if API key is provided either through env var or args
    api_key = os.getenv("LIGHTRAG_API_KEY") or args.key

    # Initialize document manager
    doc_manager = server_util.DocumentManager(args.input_dir)

    @asynccontextmanager
    async def lifespan(app: fastapi.FastAPI):
        """Lifespan context manager for startup and shutdown events"""
        # Startup logic
        pass
        yield
        # Cleanup logic (if needed)
        pass

    # Initialize FastAPI
    app = fastapi.FastAPI(
        title="LightRAG API",
        description=(
            "API for querying text using LightRAG with separate storage and input directories"
            + "(With authentication)"
            if api_key
            else ""
        ),
        version=__api_version__,
        openapi_tags=[{"name": "api"}],
        lifespan=lifespan,
    )

    # Add CORS middleware
    app.add_middleware(
        CORSMiddleware,
        allow_origins=["*"],
        allow_credentials=True,
        allow_methods=["*"],
        allow_headers=["*"],
    )

    # Create the optional API key dependency
    optional_api_key = server_util.get_api_key_dependency(api_key)

    # Create working directory if it doesn't exist
    Path(args.working_dir).mkdir(parents=True, exist_ok=True)

    async def azure_openai_model_complete(
        prompt,
        azure_ad_token,
        system_prompt=None,
        history_messages=[],
        **kwargs,
    ):
        keyword_extraction = kwargs.pop("keyword_extraction", None)
        if keyword_extraction:
            kwargs["response_format"] = GPTKeywordExtractionFormat
        if history_messages is None:
            history_messages = []
        return await az_openai.azure_openai_complete_if_cache(
            args.llm_model,
            prompt,
            system_prompt=system_prompt,
            history_messages=history_messages,
            endpoint=args.llm_binding_host,
            azure_ad_token=azure_ad_token,
            api_version=args.llm_api_version,
            **kwargs,
        )

    async def azure_embedding_function(
        texts: list[str], azure_ad_token: str
    ) -> np.ndarray:
        return await az_openai.azure_openai_embed(
            texts=texts,
            azure_ad_token=azure_ad_token,
            model=args.embedding_model,
            endpoint=args.embedding_binding_host,
            api_version=args.embedding_api_version,
        )

    async def azure_ad_token_test(azure_ad_token: str):
        try:
            await azure_openai_model_complete(
                prompt="what's the answer of 25 to the power of 3, please only output final answer.",
                azure_ad_token=azure_ad_token,
                system_prompt="You are a friendly chat bot.",
                history_messages=[],
            )
        except Exception as e:
            response = str(e)
            raise fastapi.HTTPException(status_code=401, detail=response)

    embedding_func = EmbeddingFunc(
        embedding_dim=args.embedding_dim,
        max_token_size=args.max_embed_tokens,
        func=azure_embedding_function,
    )

    # Initialize RAG
    rag = LightRAG(
        working_dir=args.working_dir,
        llm_model_func=azure_openai_model_complete,
        chunk_token_size=int(args.chunk_size),
        chunk_overlap_token_size=int(args.chunk_overlap_size),
        llm_model_kwargs={
            "timeout": args.timeout,
        },
        llm_model_name=args.llm_model,
        llm_model_max_async=args.max_async,
        llm_model_max_token_size=args.max_tokens,
        embedding_func=embedding_func,
        kv_storage=rag_storage_config.KV_STORAGE,
        graph_storage=rag_storage_config.GRAPH_STORAGE,
        vector_storage=rag_storage_config.VECTOR_STORAGE,
        doc_status_storage=rag_storage_config.DOC_STATUS_STORAGE,
        vector_db_storage_cls_kwargs={
            "cosine_better_than_threshold": args.cosine_threshold
        },
        enable_llm_cache_for_entity_extract=False,  # set to True for debuging to reduce llm fee
        embedding_cache_config={
            "enabled": True,
            "similarity_threshold": 0.95,
            "use_llm_check": False,
        },
    )

    async def index_file(file_path: Union[str, Path], azure_ad_token: str) -> None:
        """Index all files inside the folder with support for multiple file formats

        Args:
            file_path: Path to the file to be indexed (str or Path object)

        Raises:
            ValueError: If file format is not supported
            FileNotFoundError: If file doesn't exist
        """
        # Convert to Path object if string
        file_path = Path(file_path)

        # Check if file exists
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        content = ""
        # Get file extension in lowercase
        ext = file_path.suffix.lower()

        match ext:
            case ".txt" | ".md":
                # Text files handling
                async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
                    content = await f.read()

            case ".pdf" | ".docx" | ".pptx" | ".xlsx":
                async def convert_doc():
                    def sync_convert():
                        converter = DocumentConverter()
                        result = converter.convert(file_path)
                        return result.document.export_to_markdown()

                    return await asyncio.to_thread(sync_convert)

                content = await convert_doc()

            case _:
                raise ValueError(f"Unsupported file format: {ext}")

        # Insert content into RAG system
        if content:
            await rag.ainsert(content, azure_ad_token=azure_ad_token)
            doc_manager.mark_as_indexed(file_path)
            logging.info(f"Successfully indexed file: {file_path}")
        else:
            logging.warning(f"No content extracted from file: {file_path}")

    @app.post("/documents/scan", dependencies=[fastapi.Depends(optional_api_key)])
    async def scan_for_new_documents(
        background_tasks: fastapi.BackgroundTasks,
        azure_ad_token: str = fastapi.Header(None, alias="Azure_Ad_Token"),
    ):
        try:
            az_ad_token = server_util.extract_token_value(azure_ad_token)
        except:
            raise
        """Trigger the scanning process"""
        global scan_progress

        with progress_lock:
            if scan_progress["is_scanning"]:
                return {"status": "already_scanning"}

            scan_progress["is_scanning"] = True
            scan_progress["indexed_count"] = 0
            scan_progress["progress"] = 0

        # Start the scanning process in the background
        background_tasks.add_task(run_scanning_process, azure_ad_token=az_ad_token)

        return {"status": "scanning_started"}

    async def run_scanning_process(azure_ad_token: str):
        """Background task to scan and index documents"""
        global scan_progress

        try:
            new_files = doc_manager.scan_directory_for_new_files()
            scan_progress["total_files"] = len(new_files)

            for file_path in new_files:
                try:
                    with progress_lock:
                        scan_progress["current_file"] = os.path.basename(file_path)

                    await index_file(file_path, azure_ad_token=azure_ad_token)

                    with progress_lock:
                        scan_progress["indexed_count"] += 1
                        scan_progress["progress"] = (
                            scan_progress["indexed_count"]
                            / scan_progress["total_files"]
                        ) * 100

                except Exception as e:
                    logging.error(f"Error indexing file {file_path}: {str(e)}")

        except Exception as e:
            logging.error(f"Error during scanning process: {str(e)}")
        finally:
            with progress_lock:
                scan_progress["is_scanning"] = False

    @app.get("/documents/scan-progress")
    async def get_scan_progress(
        azure_ad_token: str = fastapi.Header(None, alias="Azure_Ad_Token")
    ):
        try:
            az_ad_token = server_util.extract_token_value(azure_ad_token)
        except:
            raise
        """Get the current scanning progress"""
        with progress_lock:
            return scan_progress

    @app.post("/documents/upload", dependencies=[fastapi.Depends(optional_api_key)])
    async def upload_to_input_dir(
        file: fastapi.UploadFile = fastapi.File(...),
        azure_ad_token: str = fastapi.Header(None, alias="Azure_Ad_Token"),
    ):
        """
        Endpoint for uploading a file to the input directory and indexing it.

        This API endpoint accepts a file through an HTTP POST request, checks if the
        uploaded file is of a supported type, saves it in the specified input directory,
        indexes it for retrieval, and returns a success status with relevant details.

        Parameters:
            file (UploadFile): The file to be uploaded. It must have an allowed extension as per
                               `doc_manager.supported_extensions`.

        Returns:
            dict: A dictionary containing the upload status ("success"),
                  a message detailing the operation result, and
                  the total number of indexed documents.

        Raises:
            HTTPException: If the file type is not supported, it raises a 400 Bad Request error.
                           If any other exception occurs during the file handling or indexing,
                           it raises a 500 Internal Server Error with details about the exception.
        """
        try:
            az_ad_token = server_util.extract_token_value(azure_ad_token)
        except:
            raise
        try:
            if not doc_manager.is_supported_file(file.filename):
                raise fastapi.HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type. Supported types: {doc_manager.supported_extensions}",
                )

            file_path = doc_manager.input_dir / file.filename
            with open(file_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)

            # Immediately index the uploaded file
            await index_file(file_path, azure_ad_token=az_ad_token)

            return {
                "status": "success",
                "message": f"File uploaded and indexed: {file.filename}",
                "total_documents": len(doc_manager.indexed_files),
            }
        except Exception as e:
            raise fastapi.HTTPException(status_code=500, detail=str(e))

    @app.post(
        "/query",
        response_model=server_util.QueryResponse,
        dependencies=[fastapi.Depends(optional_api_key)],
    )
    async def query_text(
        request: server_util.QueryRequest,
        azure_ad_token: str = fastapi.Header(None, alias="Azure_Ad_Token"),
    ):
        """
        Handle a POST request at the /query endpoint to process user queries using RAG capabilities.

        Parameters:
            request (QueryRequest): A Pydantic model containing the following fields:
                - query (str): The text of the user's query.
                - mode (ModeEnum): Optional. Specifies the mode of retrieval augmentation.
                - stream (bool): Optional. Determines if the response should be streamed.
                - only_need_context (bool): Optional. If true, returns only the context without further processing.

        Returns:
            QueryResponse: A Pydantic model containing the result of the query processing.
                           If a string is returned (e.g., cache hit), it's directly returned.
                           Otherwise, an async generator may be used to build the response.

        Raises:
            HTTPException: Raised when an error occurs during the request handling process,
                           with status code 500 and detail containing the exception message.
        """
        try:
            az_ad_token = server_util.extract_token_value(azure_ad_token)
        except:
            raise
        try:
            response = await rag.aquery(
                request.query,
                azure_ad_token=az_ad_token,
                param=QueryParam(
                    mode=request.mode,
                    stream=request.stream,
                    only_need_context=request.only_need_context,
                    top_k=args.top_k,
                ),
            )

            # If response is a string (e.g. cache hit), return directly
            if isinstance(response, str):
                return server_util.QueryResponse(response=response)

            # If it's an async generator, decide whether to stream based on stream parameter
            if request.stream:
                result = ""
                async for chunk in response:
                    result += chunk
                return server_util.QueryResponse(response=result)
            else:
                result = ""
                async for chunk in response:
                    result += chunk
                return server_util.QueryResponse(response=result)
        except Exception as e:
            trace_exception(e)
            raise fastapi.HTTPException(status_code=500, detail=str(e))

    @app.post("/query/stream", dependencies=[fastapi.Depends(optional_api_key)])
    async def query_text_stream(
        request: server_util.QueryRequest,
        azure_ad_token: str = fastapi.Header(None, alias="Azure_Ad_Token"),
    ):
        """
        This endpoint performs a retrieval-augmented generation (RAG) query and streams the response.

        Args:
            request (QueryRequest): The request object containing the query parameters.
            optional_api_key (Optional[str], optional): An optional API key for authentication. Defaults to None.

        Returns:
            StreamingResponse: A streaming response containing the RAG query results.
        """
        try:
            az_ad_token = server_util.extract_token_value(azure_ad_token)
        except:
            raise
        try:
            response = await rag.aquery(  # Use aquery instead of query, and add await
                request.query,
                azure_ad_token=az_ad_token,
                param=QueryParam(
                    mode=request.mode,
                    stream=True,
                    only_need_context=request.only_need_context,
                    top_k=args.top_k,
                ),
            )

            from fastapi.responses import StreamingResponse

            async def stream_generator():
                if isinstance(response, str):
                    # If it's a string, send it all at once
                    yield f"{json.dumps({'response': response})}\n"
                else:
                    # If it's an async generator, send chunks one by one
                    try:
                        async for chunk in response:
                            if chunk:  # Only send non-empty content
                                yield f"{json.dumps({'response': chunk})}\n"
                    except Exception as e:
                        logging.error(f"Streaming error: {str(e)}")
                        yield f"{json.dumps({'error': str(e)})}\n"

            return StreamingResponse(
                stream_generator(),
                media_type="application/x-ndjson",
                headers={
                    "Cache-Control": "no-cache",
                    "Connection": "keep-alive",
                    "Content-Type": "application/x-ndjson",
                    "Access-Control-Allow-Origin": "*",
                    "Access-Control-Allow-Methods": "POST, OPTIONS",
                    "Access-Control-Allow-Headers": "Content-Type",
                    "X-Accel-Buffering": "no",  # Disable Nginx buffering
                },
            )
        except Exception as e:
            trace_exception(e)
            raise fastapi.HTTPException(status_code=500, detail=str(e))

    @app.post(
        "/documents/text",
        response_model=server_util.InsertResponse,
        dependencies=[fastapi.Depends(optional_api_key)],
    )
    async def insert_text(
        request: server_util.InsertTextRequest,
        azure_ad_token: str = fastapi.Header(None, alias="Azure_Ad_Token"),
    ):
        """
        Insert text into the Retrieval-Augmented Generation (RAG) system.

        This endpoint allows you to insert text data into the RAG system for later retrieval and use in generating responses.

        Args:
            request (InsertTextRequest): The request body containing the text to be inserted.

        Returns:
            InsertResponse: A response object containing the status of the operation, a message, and the number of documents inserted.
        """
        try:
            az_ad_token = server_util.extract_token_value(azure_ad_token)
        except:
            raise
        try:
            await rag.ainsert(request.text, azure_ad_token=az_ad_token)
            return server_util.InsertResponse(
                status="success",
                message="Text successfully inserted",
                document_count=1,
            )
        except Exception as e:
            raise fastapi.HTTPException(status_code=500, detail=str(e))

    @app.post(
        "/documents/file",
        response_model=server_util.InsertResponse,
        dependencies=[fastapi.Depends(optional_api_key)],
    )
    async def insert_file(
        file: fastapi.UploadFile = fastapi.File(...),
        description: str = fastapi.Form(None),
        azure_ad_token: str = fastapi.Header(None, alias="Azure_Ad_Token"),
    ):
        """Insert a file directly into the RAG system

        Args:
            file: Uploaded file
            description: Optional description of the file

        Returns:
            InsertResponse: Status of the insertion operation

        Raises:
            HTTPException: For unsupported file types or processing errors
        """
        try:
            az_ad_token = server_util.extract_token_value(azure_ad_token)
        except:
            raise
        try:
            content = ""
            # Get file extension in lowercase
            ext = Path(file.filename).suffix.lower()

            match ext:
                case ".txt" | ".md":
                    # Text files handling
                    text_content = await file.read()
                    content = text_content.decode("utf-8")

                case ".pdf" | ".docx" | ".pptx" | ".xlsx":
                    from docling.document_converter import DocumentConverter

                    # Create a temporary file to save the uploaded content
                    temp_path = Path("temp") / file.filename
                    temp_path.parent.mkdir(exist_ok=True)

                    # Save the uploaded file
                    with temp_path.open("wb") as f:
                        f.write(await file.read())

                    try:
                        converter = DocumentConverter()
                        result = converter.convert(str(temp_path))
                        content = result.document.export_to_markdown()
                    finally:
                        # Clean up the temporary file
                        temp_path.unlink()

            # Insert content into RAG system
            if content:
                # Add description if provided
                if description:
                    content = f"{description}\n\n{content}"

                await rag.ainsert(content, az_ad_token)
                logging.info(f"Successfully indexed file: {file.filename}")

                return server_util.InsertResponse(
                    status="success",
                    message=f"File '{file.filename}' successfully inserted",
                    document_count=1,
                )
            else:
                raise fastapi.HTTPException(
                    status_code=400,
                    detail="No content could be extracted from the file",
                )

        except UnicodeDecodeError:
            raise fastapi.HTTPException(
                status_code=400, detail="File encoding not supported"
            )
        except Exception as e:
            logging.error(f"Error processing file {file.filename}: {str(e)}")
            raise fastapi.HTTPException(status_code=500, detail=str(e))

    @app.post(
        "/documents/batch",
        response_model=server_util.InsertResponse,
        dependencies=[fastapi.Depends(optional_api_key)],
    )
    async def insert_batch(
        files: List[fastapi.UploadFile] = fastapi.File(...),
        azure_ad_token: str = fastapi.Header(None, alias="Azure_Ad_Token"),
    ):
        """Process multiple files in batch mode

        Args:
            files: List of files to process

        Returns:
            InsertResponse: Status of the batch insertion operation

        Raises:
            HTTPException: For processing errors
        """
        try:
            az_ad_token = server_util.extract_token_value(azure_ad_token)
        except:
            raise
        try:
            inserted_count = 0
            failed_files = []

            for file in files:
                try:
                    content = ""
                    ext = Path(file.filename).suffix.lower()

                    match ext:
                        case ".txt" | ".md":
                            text_content = await file.read()
                            content = text_content.decode("utf-8")

                        case ".pdf":
                            pdf_content = await file.read()
                            pdf_file = BytesIO(pdf_content)
                            reader = PdfReader(pdf_file)
                            for page in reader.pages:
                                content += page.extract_text() + "\n"

                        case ".docx":

                            docx_content = await file.read()
                            docx_file = BytesIO(docx_content)
                            doc = Document(docx_file)
                            content = "\n".join(
                                [paragraph.text for paragraph in doc.paragraphs]
                            )

                        case ".pptx":
                            pptx_content = await file.read()
                            pptx_file = BytesIO(pptx_content)
                            prs = Presentation(pptx_file)
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        content += shape.text + "\n"

                        case _:
                            failed_files.append(f"{file.filename} (unsupported type)")
                            continue

                    if content:
                        await rag.ainsert(content, azure_ad_token=az_ad_token)
                        inserted_count += 1
                        logging.info(f"Successfully indexed file: {file.filename}")
                    else:
                        failed_files.append(f"{file.filename} (no content extracted)")

                except UnicodeDecodeError:
                    failed_files.append(f"{file.filename} (encoding error)")
                except Exception as e:
                    failed_files.append(f"{file.filename} ({str(e)})")
                    logging.error(f"Error processing file {file.filename}: {str(e)}")

            # Prepare status message
            if inserted_count == len(files):
                status = "success"
                status_message = f"Successfully inserted all {inserted_count} documents"
            elif inserted_count > 0:
                status = "partial_success"
                status_message = f"Successfully inserted {inserted_count} out of {len(files)} documents"
                if failed_files:
                    status_message += f". Failed files: {', '.join(failed_files)}"
            else:
                status = "failure"
                status_message = "No documents were successfully inserted"
                if failed_files:
                    status_message += f". Failed files: {', '.join(failed_files)}"

            return server_util.InsertResponse(
                status=status,
                message=status_message,
                document_count=inserted_count,
            )

        except Exception as e:
            logging.error(f"Batch processing error: {str(e)}")
            raise fastapi.HTTPException(status_code=500, detail=str(e))

    @app.delete(
        "/documents",
        response_model=server_util.InsertResponse,
        dependencies=[fastapi.Depends(optional_api_key)],
    )
    async def clear_documents(
        azure_ad_token: str = fastapi.Header(None, alias="Azure_Ad_Token")
    ):
        """
        Clear all documents from the LightRAG system.

        This endpoint deletes all text chunks, entities vector database, and relationships vector database,
        effectively clearing all documents from the LightRAG system.

        Returns:
            InsertResponse: A response object containing the status, message, and the new document count (0 in this case).
        """
        try:
            az_ad_token = server_util.extract_token_value(azure_ad_token)
        except:
            raise
        try:
            rag.text_chunks = []
            rag.entities_vdb = None
            rag.relationships_vdb = None
            return server_util.InsertResponse(
                status="success",
                message="All documents cleared successfully",
                document_count=0,
            )
        except Exception as e:
            raise fastapi.HTTPException(status_code=500, detail=str(e))

    # query all graph labels
    @app.get("/graph/label/list")
    async def get_graph_labels(
        azure_ad_token: str = fastapi.Header(None, alias="Azure_Ad_Token")
    ):
        return await rag.get_graph_labels()

    # query all graph
    @app.get("/graphs")
    async def get_graphs(
        label: str, azure_ad_token: str = fastapi.Header(None, alias="Azure_Ad_Token")
    ):
        return await rag.get_graps(nodel_label=label, max_depth=100)

    @app.get("/documents", dependencies=[fastapi.Depends(optional_api_key)])
    async def documents(
        azure_ad_token: str = fastapi.Header(None, alias="Azure_Ad_Token")
    ):
        """Get current system status"""
        return doc_manager.indexed_files

    @app.get("/health", dependencies=[fastapi.Depends(optional_api_key)])
    async def get_status(authorization: Optional[str] = fastapi.Header(None)):
        try:
            az_ad_token = server_util.extract_token_value(authorization)
            await azure_ad_token_test(az_ad_token)
        except:
            raise
        value_to_check = [
            "llm_binding_host",
            "llm_model",
            "llm_api_version",
            "embedding_binding_host",
            "embedding_model",
            "embedding_api_version",
        ]
        unhealthy = any(getattr(args, prop) is None for prop in value_to_check)
        """Get current system status"""
        files = doc_manager.scan_directory()
        return {
            "Status": "Healthy" if not unhealthy else "Unhealthy",
            "Working Directory": str(args.working_dir),
            "Input Directory": str(args.input_dir),
            "Indexed Files": [str(f) for f in files],
            "Indexed Files Count": len(files),
            "Configuration": {
                "Azure Open AI Endpoint": args.llm_binding_host,
                "Azure Open AI Large Language Model Name": args.llm_model,
                "Azure Open AI API Version": args.llm_api_version,
                "Azure Open AI Embedding Endpoint": args.embedding_binding_host,
                "Azure Open AI Embedding Model Name": args.embedding_model,
                "Azure Open AI Embedding API Version": args.embedding_api_version,
                "Max Tokens": args.max_tokens,
                "KV Storage": rag_storage_config.KV_STORAGE,
                "Doc Status Storage": rag_storage_config.DOC_STATUS_STORAGE,
                "Graph Storage": rag_storage_config.GRAPH_STORAGE,
                "Vector Storage": rag_storage_config.VECTOR_STORAGE,
            },
        }

    # Serve the static files
    static_dir = Path(__file__).parent / "static"
    static_dir.mkdir(exist_ok=True)
    app.mount("/", StaticFiles(directory=static_dir, html=True), name="static")

    return app
