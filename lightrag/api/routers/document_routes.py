"""
This module contains all document-related routes for the LightRAG API.
"""

import asyncio
import io
import logging
import os
import traceback
import pipmaster as pm
from azure.storage.blob import BlobServiceClient, BlobLeaseClient, ContainerClient
from ..rag_instance_manager import RAGInstanceManager
from typing import Dict, List, Optional, Any
from fastapi import (
    APIRouter,
    Depends,
    File,
    HTTPException,
    UploadFile,
    Header,
)
from fastapi.responses import JSONResponse
from pydantic import BaseModel, Field, field_validator
from ... import LightRAG
from ...az_token_credential import LightRagTokenCredential
from ...base import DocProcessingStatus, DocStatus, StoragesStatus
from ..utils_api import (
    get_api_key_dependency,
    initialize_rag_with_header,
    wait_for_storage_initialization,
    get_lightrag_token_credential,
    extract_token_value,
    try_get_container_lease,
)


router = APIRouter(prefix="/documents", tags=["documents"])

# Lock for thread-safe operations
progress_lock = asyncio.Lock()

# Temporary file prefix
temp_prefix = "__tmp__"


class InsertTextRequest(BaseModel):
    text: str = Field(
        min_length=1,
        description="The text to insert",
    )

    id: Optional[str] = Field(
        default=None, min_length=1, description="The id of the text to insert"
    )

    @field_validator("text", mode="after")
    @classmethod
    def strip_after(cls, text: str) -> str:
        return text.strip()


class InsertTextsRequest(BaseModel):
    texts: list[str] = Field(
        min_length=1,
        description="The texts to insert",
    )

    ids: Optional[list[str]] = Field(
        default=None, min_length=1, description="The ids of the texts to insert"
    )

    @field_validator("texts", mode="after")
    @classmethod
    def strip_after(cls, texts: list[str]) -> list[str]:
        return [text.strip() for text in texts]


class InsertResponse(BaseModel):
    status: str = Field(description="Status of the operation")
    message: str = Field(description="Message describing the operation result")


class DocStatusResponse(BaseModel):
    @staticmethod
    def format_datetime(dt: Any) -> Optional[str]:
        if dt is None:
            return None
        if isinstance(dt, str):
            return dt
        return dt.isoformat()

    """Response model for document status

    Attributes:
        id: Document identifier
        content_summary: Summary of document content
        content_length: Length of document content
        status: Current processing status
        created_at: Creation timestamp (ISO format string)
        updated_at: Last update timestamp (ISO format string)
        chunks_count: Number of chunks (optional)
        error: Error message if any (optional)
        metadata: Additional metadata (optional)
    """

    id: str
    content_summary: str
    content_length: int
    status: DocStatus
    created_at: str
    updated_at: str
    chunks_count: Optional[int] = None
    error: Optional[str] = None
    metadata: Optional[dict[str, Any]] = None


class DocsStatusesResponse(BaseModel):
    statuses: Dict[DocStatus, List[DocStatusResponse]] = {}


async def get_file_stream_from_storage(
    file_name: str, container_client: ContainerClient
):
    blob_lease = None
    lease = None
    try:
        lease: BlobLeaseClient = await try_get_container_lease(container_client)
        stream = io.BytesIO()
        blob_client = container_client.get_blob_client(file_name)
        blob_lease: BlobLeaseClient = await try_get_container_lease(blob_client)
        blob_client.download_blob().readinto(stream)
        stream.seek(0, os.SEEK_SET)
        return stream
    except Exception as ex:
        logging.error(f"Error getting file stream from storage: {str(ex)}")
        raise HTTPException(status_code=500, detail=str(ex))
    finally:
        if blob_lease:
            blob_lease.release()
        if lease:
            lease.release()


async def pipeline_enqueue_file(
    rag: LightRAG,
    file_path: str,
    storage_account_url: str,
    storage_container_name: str,
    access_token: LightRagTokenCredential,
) -> bool:
    """Add a file to the queue for processing

    Args:
        rag: LightRAG instance
        file_path: Path to the saved file
    Returns:
        bool: True if the file was successfully enqueued, False otherwise
    """

    try:
        content = ""
        root, ext = os.path.splitext(file_path)

        blob_client = BlobServiceClient(
            account_url=storage_account_url, credential=access_token
        )
        container_client = blob_client.get_container_client(storage_container_name)
        container_client.get_container_properties()  # this is to check if the container exists and authentication is valid
        file_name = f"{rag.document_manager.input_dir}/{file_path}"
        lease: BlobLeaseClient = await try_get_container_lease(container_client)
        blob_list = container_client.list_blob_names()
        if file_name not in blob_list:
            logging.error(f"File {file_name} not found in storage")
            lease.release()
            raise HTTPException(
                status_code=404,
                detail=f"File {file_name} not found in storage {storage_container_name}",
            )
        lease.release()

        # Process based on file type
        match ext:
            case (
                ".txt"
                | ".md"
                | ".html"
                | ".htm"
                | ".tex"
                | ".json"
                | ".xml"
                | ".yaml"
                | ".yml"
                | ".rtf"
                | ".odt"
                | ".epub"
                | ".csv"
                | ".log"
                | ".conf"
                | ".ini"
                | ".properties"
                | ".sql"
                | ".bat"
                | ".sh"
                | ".c"
                | ".cpp"
                | ".py"
                | ".java"
                | ".js"
                | ".ts"
                | ".swift"
                | ".go"
                | ".rb"
                | ".php"
                | ".css"
                | ".scss"
                | ".less"
            ):
                lease = None
                blob_lease = None
                try:
                    lease: BlobLeaseClient = await try_get_container_lease(
                        container_client
                    )
                    blob_client = container_client.get_blob_client(file_name)
                    blob_lease: BlobLeaseClient = await try_get_container_lease(
                        blob_client
                    )
                    file_byte = container_client.download_blob(file_name).readall()
                    content = file_byte.decode("utf-8")
                except Exception as ex:
                    logging.error(
                        f"Error processing or enqueueing file {file_path}: {str(ex)}"
                    )
                    raise HTTPException(status_code=500, detail=str(ex))
                finally:
                    if blob_lease:
                        blob_lease.release()
                    if lease:
                        lease.release()

            case ".pdf":
                if not pm.is_installed("pypdf2"):
                    pm.install("pypdf2")
                from PyPDF2 import PdfReader  # type: ignore

                pdf_file = await get_file_stream_from_storage(
                    file_name, container_client
                )
                reader = PdfReader(pdf_file)
                for page in reader.pages:
                    content += page.extract_text() + "\n"
            case ".docx":
                if not pm.is_installed("python-docx"):
                    pm.install("python-docx")
                from docx import Document
                from io import BytesIO

                docx_file = await get_file_stream_from_storage(
                    file_name, container_client
                )
                doc = Document(docx_file)
                content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            case ".pptx":
                if not pm.is_installed("python-pptx"):
                    pm.install("python-pptx")
                from pptx import Presentation
                from io import BytesIO

                pptx_file = await get_file_stream_from_storage(
                    file_name, container_client
                )
                prs = Presentation(pptx_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
            case ".xlsx":
                if not pm.is_installed("openpyxl"):
                    pm.install("openpyxl")
                from openpyxl import load_workbook
                from io import BytesIO

                xlsx_file = await get_file_stream_from_storage(
                    file_name, container_client
                )
                wb = load_workbook(xlsx_file)
                for sheet in wb:
                    content += f"Sheet: {sheet.title}\n"
                    for row in sheet.iter_rows(values_only=True):
                        content += (
                            "\t".join(
                                str(cell) if cell is not None else "" for cell in row
                            )
                            + "\n"
                        )
                    content += "\n"
            case _:
                logging.error(f"Unsupported file type: {file_path} (extension {ext})")
                return False

        # Insert into the RAG queue
        if content:
            await rag.apipeline_enqueue_documents(
                storage_account_url, storage_container_name, access_token, content
            )
            logging.info(f"Successfully fetched and enqueued file: {file_path}")
            return True
        else:
            logging.error(f"No content could be extracted from file: {file_path}")

    except Exception as e:
        logging.error(f"Error processing or enqueueing file {file_path}: {str(e)}")
        logging.error(traceback.format_exc())
    return False


async def pipeline_index_file(
    rag: LightRAG,
    file_path: str,
    ai_access_token: str,
    storage_account_url: str,
    storage_container_name: str,
    access_token: LightRagTokenCredential,
):
    """Index a file"""
    try:
        if await pipeline_enqueue_file(
            rag, file_path, storage_account_url, storage_container_name, access_token
        ):
            await rag.apipeline_process_enqueue_documents(
                ai_access_token,
                storage_account_url,
                storage_container_name,
                access_token,
            )

    except Exception as e:
        logging.error(f"Error indexing file {file_path}: {str(e)}")
        logging.error(traceback.format_exc())


async def pipeline_index_texts(
    rag: LightRAG,
    texts: List[str],
    ids: List[str] | None,
    ai_access_token: str,
    storage_account_url: str,
    storage_container_name: str,
    access_token: LightRagTokenCredential,
):
    """Index a list of texts"""
    if not texts:
        return
    await rag.apipeline_enqueue_documents(
        storage_account_url, storage_container_name, access_token, texts, ids
    )
    await rag.apipeline_process_enqueue_documents(
        ai_access_token, storage_account_url, storage_container_name, access_token
    )


async def empty_light_rag_databases(
    rag: LightRAG,
    storage_account_url: str,
    storage_container_name: str,
    access_token: LightRagTokenCredential,
):
    blob_client = BlobServiceClient(
        account_url=storage_account_url, credential=access_token
    )
    container_client = blob_client.get_container_client(storage_container_name)
    container_client.get_container_properties()  # this is to check if the container exists and authentication is valid
    blobs_list = container_client.list_blobs(name_starts_with=rag.working_dir)
    for blob in blobs_list:
        blob_client = container_client.get_blob_client(blob)
        blob_lease = await try_get_container_lease(blob_client)
        blob_client.delete_blob(lease=blob_lease)
    await rag.initialize_storages(access_token)


async def upload_file(
    storage_account_url: str,
    storage_container_name: str,
    access_token: LightRagTokenCredential,
    rag: LightRAG,
    file: UploadFile = File(...),
) -> str:
    """Save the uploaded file to a temporary location"""
    blob_client = BlobServiceClient(
        account_url=storage_account_url, credential=access_token
    )
    container_client = blob_client.get_container_client(storage_container_name)
    container_client.get_container_properties()  # this is to check if the container exists and authentication is valid
    file_name = f"{rag.document_manager.input_dir}/{file.filename}"
    lease = None
    while lease is None:
        try:
            lease: BlobLeaseClient = container_client.acquire_lease()
            break
        except Exception as e:
            logging.error(f"Error acquiring lease: {str(e)}")
            await asyncio.sleep(5)
    blob_list = container_client.list_blob_names()
    if file_name in blob_list:
        logging.error(f"File {file_name} already exists in storage")
        raise HTTPException(
            status_code=400,
            detail=f"File {file_name} already exists in storage {storage_container_name}",
        )
    else:
        blob_name = file_name
        blob_client = container_client.get_blob_client(blob_name)
        file_bytes = await file.read()
        blob_client.upload_blob(file_bytes, overwrite=False)
    lease.release()
    return file.filename


async def run_scanning_process(
    rag: LightRAG,
    ai_access_token: str,
    storage_account_url: str,
    storage_container_name: str,
    access_token: LightRagTokenCredential,
):
    """Background task to scan and index documents"""
    try:
        new_files = rag.document_manager.scan_directory_for_new_file(
            storage_account_url, storage_container_name, access_token
        )
        if new_files is None:
            return
        rag.scan_progress["total_files"] = len(new_files)

        logging.info(f"Found {len(new_files)} new files to index.")
        for file_path in new_files:
            try:
                async with progress_lock:
                    rag.scan_progress["current_file"] = new_files
                await pipeline_index_file(
                    rag,
                    file_path,
                    ai_access_token,
                    storage_account_url,
                    storage_container_name,
                    access_token,
                )

                async with progress_lock:
                    rag.scan_progress["indexed_count"] += 1
                    rag.scan_progress["progress"] = (
                        rag.scan_progress["indexed_count"]
                        / rag.scan_progress["total_files"]
                    ) * 100

            except Exception as e:
                logging.error(f"Error indexing file {file_path}: {str(e)}")

    except Exception as e:
        logging.error(f"Error during scanning process: {str(e)}")
    finally:
        async with progress_lock:
            rag.scan_progress["is_scanning"] = False


def create_document_routes(
    rag_instance_manager: RAGInstanceManager, api_key: Optional[str] = None
):
    optional_api_key = get_api_key_dependency(api_key)

    @router.post(
        "/text", response_model=InsertResponse, dependencies=[Depends(optional_api_key)]
    )
    async def insert_text(
        request: InsertTextRequest,
        storage_account_url: str = Header(alias="Storage_Account_Url"),
        storage_container_name: str = Header(alias="Storage_Container_Name"),
        storage_token_expiry: str = Header(
            default=None, alias="Storage_Access_Token_Expiry"
        ),
        ai_access_token: str = Header(alias="Azure-AI-Access-Token"),
        storage_access_token: str = Header(alias="Storage_Access_Token"),
        X_Affinity_Token: str = Header(None, alias="X-Affinity-Token"),
    ) -> JSONResponse:
        """
        Insert text into the RAG system.

        This endpoint allows you to insert text data into the RAG system for later retrieval
        and use in generating responses.
        """
        try:
            ai_access_token = extract_token_value(
                ai_access_token, "Azure-AI-Access-Token"
            )
            storage_access_token = extract_token_value(
                storage_access_token, "Storage_Access_Token"
            )
            rag = initialize_rag_with_header(
                rag_instance_manager,
                storage_account_url,
                storage_container_name,
                X_Affinity_Token,
                storage_access_token,
                storage_token_expiry,
            )
            await wait_for_storage_initialization(
                rag,
                get_lightrag_token_credential(
                    storage_access_token, storage_token_expiry
                ),
            )
            await pipeline_index_texts(
                rag,
                [request.text],
                [request.id] if request.id else None,
                ai_access_token,
                storage_account_url,
                storage_container_name,
                get_lightrag_token_credential(
                    storage_access_token, storage_token_expiry
                ),
            )
            response = JSONResponse(
                content={
                    "status": "success",
                    "message": "Text successfully received and indexed.",
                },
                headers={"X-Affinity-Token": rag.affinity_token},
            )
            return response
        except Exception as e:
            logging.error(f"Error /documents/text: {str(e)}")
            logging.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/texts",
        response_model=InsertResponse,
        dependencies=[Depends(optional_api_key)],
    )
    async def insert_texts(
        request: InsertTextsRequest,
        storage_account_url: str = Header(alias="Storage_Account_Url"),
        storage_container_name: str = Header(alias="Storage_Container_Name"),
        storage_token_expiry: str = Header(
            default=None, alias="Storage_Access_Token_Expiry"
        ),
        ai_access_token: str = Header(alias="Azure-AI-Access-Token"),
        storage_access_token: str = Header(alias="Storage_Access_Token"),
        X_Affinity_Token: str = Header(None, alias="X-Affinity-Token"),
    ):
        """
        Insert multiple texts into the RAG system.

        This endpoint allows you to insert multiple text entries into the RAG system
        in a single request.
        """
        try:
            ai_access_token = extract_token_value(
                ai_access_token, "Azure-AI-Access-Token"
            )
            storage_access_token = extract_token_value(
                storage_access_token, "Storage_Access_Token"
            )
            rag = initialize_rag_with_header(
                rag_instance_manager,
                storage_account_url,
                storage_container_name,
                X_Affinity_Token,
                storage_access_token,
                storage_token_expiry,
            )
            await wait_for_storage_initialization(
                rag,
                get_lightrag_token_credential(
                    storage_access_token, storage_token_expiry
                ),
            )
            await pipeline_index_texts(
                rag,
                request.texts,
                request.ids if request.ids else None,
                ai_access_token,
                storage_account_url,
                storage_container_name,
                get_lightrag_token_credential(
                    storage_access_token, storage_token_expiry
                ),
            )
            response = JSONResponse(
                content={
                    "status": "success",
                    "message": "Text successfully received and indexed",
                },
                headers={"X-Affinity-Token": rag.affinity_token},
            )
            return response
        except Exception as e:
            logging.error(f"Error /documents/text: {str(e)}")
            logging.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/file", response_model=InsertResponse, dependencies=[Depends(optional_api_key)]
    )
    async def insert_file(
        file: UploadFile = File(...),
        storage_account_url: str = Header(alias="Storage_Account_Url"),
        storage_container_name: str = Header(alias="Storage_Container_Name"),
        storage_token_expiry: str = Header(
            default=None, alias="Storage_Access_Token_Expiry"
        ),
        ai_access_token: str = Header(alias="Azure-AI-Access-Token"),
        storage_access_token: str = Header(alias="Storage_Access_Token"),
        X_Affinity_Token: str = Header(None, alias="X-Affinity-Token"),
    ):
        """
        Insert a file directly into the RAG system.

        This endpoint accepts a file upload and processes it for inclusion in the RAG system.
        The file is saved temporarily and processed in the background.
        """

        try:
            ai_access_token = extract_token_value(
                ai_access_token, "Azure-AI-Access-Token"
            )
            storage_access_token = extract_token_value(
                storage_access_token, "Storage_Access_Token"
            )
            rag = initialize_rag_with_header(
                rag_instance_manager,
                storage_account_url,
                storage_container_name,
                X_Affinity_Token,
                storage_access_token,
                storage_token_expiry,
            )
            if not rag.document_manager.is_supported_file(file.filename):
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type. Supported types: {rag.document_manager.supported_extensions}",
                )
            await wait_for_storage_initialization(
                rag,
                get_lightrag_token_credential(
                    storage_access_token, storage_token_expiry
                ),
            )
            temp_path = await upload_file(
                storage_account_url,
                storage_container_name,
                get_lightrag_token_credential(
                    storage_access_token, storage_token_expiry
                ),
                rag,
                file,
            )
            await pipeline_index_file(
                rag,
                temp_path,
                ai_access_token,
                storage_account_url,
                storage_container_name,
                get_lightrag_token_credential(
                    storage_access_token, storage_token_expiry
                ),
            )
            response = JSONResponse(
                content={
                    "status": "success",
                    "message": f"File '{file.filename}' uploaded and indexed successfully.",
                },
                headers={"X-Affinity-Token": rag.affinity_token},
            )
            return response
        except Exception as e:
            logging.error(f"Error /documents/file: {str(e)}")
            logging.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.delete(
        "", response_model=InsertResponse, dependencies=[Depends(optional_api_key)]
    )
    async def clear_documents(
        storage_account_url: str = Header(alias="Storage_Account_Url"),
        storage_container_name: str = Header(alias="Storage_Container_Name"),
        storage_token_expiry: str = Header(
            default=None, alias="Storage_Access_Token_Expiry"
        ),
        ai_access_token: str = Header(alias="Azure-AI-Access-Token"),
        storage_access_token: str = Header(alias="Storage_Access_Token"),
        X_Affinity_Token: str = Header(None, alias="X-Affinity-Token"),
    ):
        """
        Clear all documents from the RAG system.

        This endpoint deletes all text chunks, entities vector database, and relationships
        vector database, effectively clearing all documents from the RAG system.
        """
        try:
            ai_access_token = extract_token_value(
                ai_access_token, "Azure-AI-Access-Token"
            )
            storage_access_token = extract_token_value(
                storage_access_token, "Storage_Access_Token"
            )
            rag = initialize_rag_with_header(
                rag_instance_manager,
                storage_account_url,
                storage_container_name,
                X_Affinity_Token,
                storage_access_token,
                storage_token_expiry,
            )
            rag.text_chunks = []
            rag.entities_vdb = None
            rag.relationships_vdb = None
            await empty_light_rag_databases(
                rag,
                storage_account_url,
                storage_container_name,
                get_lightrag_token_credential(
                    storage_access_token, storage_token_expiry
                ),
            )
            response = JSONResponse(
                content={
                    "status": "success",
                    "message": "All databases has been deleted in remote storage, new empty databases has been created.",
                },
                headers={"X-Affinity-Token": rag.affinity_token},
            )
            return response
        except Exception as e:
            logging.error(f"Error DELETE /documents: {str(e)}")
            logging.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.get("", dependencies=[Depends(optional_api_key)])
    async def documents(
        storage_account_url: str = Header(alias="Storage_Account_Url"),
        storage_container_name: str = Header(alias="Storage_Container_Name"),
        storage_token_expiry: str = Header(
            default=None, alias="Storage_Access_Token_Expiry"
        ),
        ai_access_token: str = Header(alias="Azure-AI-Access-Token"),
        storage_access_token: str = Header(alias="Storage_Access_Token"),
        X_Affinity_Token: str = Header(None, alias="X-Affinity-Token"),
    ) -> JSONResponse:
        """
        Get the status of all documents in the system.

        This endpoint retrieves the current status of all documents, grouped by their
        processing status (PENDING, PROCESSING, PROCESSED, FAILED).
        """

        try:
            ai_access_token = extract_token_value(
                ai_access_token, "Azure-AI-Access-Token"
            )
            storage_access_token = extract_token_value(
                storage_access_token, "Storage_Access_Token"
            )
            statuses = (
                DocStatus.PENDING,
                DocStatus.PROCESSING,
                DocStatus.PROCESSED,
                DocStatus.FAILED,
            )
            rag = initialize_rag_with_header(
                rag_instance_manager,
                storage_account_url,
                storage_container_name,
                X_Affinity_Token,
                storage_access_token,
                storage_token_expiry,
            )
            await wait_for_storage_initialization(
                rag,
                get_lightrag_token_credential(
                    storage_access_token, storage_token_expiry
                ),
            )
            tasks = [rag.get_docs_by_status(status) for status in statuses]
            results: List[Dict[str, DocProcessingStatus]] = await asyncio.gather(*tasks)

            response = DocsStatusesResponse()
            for idx, result in enumerate(results):
                status = statuses[idx]
                for doc_id, doc_status in result.items():
                    if status not in response.statuses:
                        response.statuses[status] = []
                    response.statuses[status].append(
                        DocStatusResponse(
                            id=doc_id,
                            content_summary=doc_status.content_summary,
                            content_length=doc_status.content_length,
                            status=doc_status.status,
                            created_at=DocStatusResponse.format_datetime(
                                doc_status.created_at
                            ),
                            updated_at=DocStatusResponse.format_datetime(
                                doc_status.updated_at
                            ),
                            chunks_count=doc_status.chunks_count,
                            error=doc_status.error,
                            metadata=doc_status.metadata,
                        )
                    )
            return JSONResponse(
                content=response.model_dump(),
                headers={"X-Affinity-Token": rag.affinity_token},
            )
        except Exception as e:
            logging.error(f"Error GET /documents: {str(e)}")
            logging.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    return router
